"""
PowerPoint slide generation using python-pptx
"""

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pathlib import Path
from typing import List, Dict

from .error_handler import ErrorHandler


class SlideGenerator:
    """Generate PowerPoint presentation from allocation plan"""
    
    def __init__(self, template_path: str):
        self.template_path = Path(template_path)
        self.prs: Presentation = None
        self.error_handler = ErrorHandler()
        
        # Validate template exists
        ErrorHandler.validate_file_exists(str(self.template_path))
    
    def generate(
        self, 
        allocation_plan: List[Dict], 
        output_path: str
    ) -> str:
        """
        Generate final PowerPoint presentation
        
        Args:
            allocation_plan: List of slide allocations
            output_path: Path for output file
            
        Returns:
            Path to generated file
        """
        self.error_handler.log_info(
            f"Generating presentation: {len(allocation_plan)} slides"
        )
        
        # Load fresh template
        self.prs = Presentation(str(self.template_path))
        
        # Clear existing slides
        self._clear_all_slides()
        
        # Generate slides from plan
        for idx, slide_plan in enumerate(allocation_plan):
            try:
                self._create_slide(slide_plan)
            except Exception as e:
                self.error_handler.log_error(
                    f"Error creating slide {idx + 1}: {e}", e
                )
                # Continue with next slide
                continue
        
        # Save output
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        self.prs.save(str(output_path))
        
        self.error_handler.log_info(f"✅ Presentation saved: {output_path}")
        
        return str(output_path)
    
    def _clear_all_slides(self):
        """Remove all slides from presentation"""
        xml_slides = self.prs.slides._sldIdLst
        slides_to_remove = list(range(len(xml_slides)))
        
        for idx in reversed(slides_to_remove):
            xml_slides.remove(xml_slides[idx])
    
    def _create_slide(self, slide_plan: Dict):
        """
        Create individual slide from plan
        
        Args:
            slide_plan: Allocation dictionary for this slide
        """
        # Get layout index
        layout_idx = slide_plan.get('template_slide_idx', 0)
        
        # Ensure layout index is valid
        if layout_idx >= len(self.prs.slide_layouts):
            self.error_handler.log_warning(
                f"Invalid layout index {layout_idx}, using layout 0"
            )
            layout_idx = 0
        
        # Add slide with layout
        layout = self.prs.slide_layouts[layout_idx]
        slide = self.prs.slides.add_slide(layout)
        
        # Fill placeholders
        content_mapping = slide_plan.get('content', {})
        
        for ph_idx, content_data in content_mapping.items():
            self._fill_placeholder(slide, ph_idx, content_data)
        
        # Add notes if present
        notes = slide_plan.get('notes', '')
        if notes:
            self._add_notes(slide, notes)
    
    def _fill_placeholder(
        self, 
        slide, 
        ph_idx: int, 
        content_data: Dict
    ):
        """
        Fill specific placeholder with content
        
        Args:
            slide: Slide object
            ph_idx: Placeholder index
            content_data: Content dictionary
        """
        try:
            placeholder = slide.placeholders[ph_idx]
        except KeyError:
            self.error_handler.log_warning(
                f"Placeholder {ph_idx} not found, attempting fallback"
            )
            self._add_text_box_fallback(slide, content_data)
            return
        
        content_type = content_data.get('type', 'plain')
        content_format = content_data.get('format', 'plain')
        
        if content_type == 'title':
            self._fill_title(placeholder, content_data)
        elif content_type == 'body':
            if content_format == 'bullets':
                self._fill_bullets(placeholder, content_data)
            else:
                self._fill_plain_text(placeholder, content_data)
        else:
            self._fill_plain_text(placeholder, content_data)
    
    def _fill_title(self, placeholder, content_data: Dict):
        """Fill title placeholder"""
        text = content_data.get('text', '')
        placeholder.text = text
    
    def _fill_bullets(self, placeholder, content_data: Dict):
        """Fill placeholder with bullet points"""
        items = content_data.get('text', [])
        
        if not items:
            return
        
        # Ensure items is a list
        if isinstance(items, str):
            items = [items]
        
        text_frame = placeholder.text_frame
        text_frame.clear()  # Clear existing content
        
        # Add bullet points
        for idx, bullet in enumerate(items):
            if idx == 0:
                # Use first paragraph
                p = text_frame.paragraphs[0]
            else:
                # Add new paragraph
                p = text_frame.add_paragraph()
            
            p.text = str(bullet)
            p.level = 0  # Top-level bullet
            
            # Optional: Set bullet formatting
            # p.font.size = Pt(18)
    
    def _fill_plain_text(self, placeholder, content_data: Dict):
        """Fill placeholder with plain text"""
        text = content_data.get('text', '')
        
        if isinstance(text, list):
            text = '\n'.join(str(item) for item in text)
        
        placeholder.text = str(text)
    
    def _add_text_box_fallback(self, slide, content_data: Dict):
        """
        Add text box as fallback when placeholder not found
        
        Args:
            slide: Slide object
            content_data: Content dictionary
        """
        self.error_handler.log_info("Using text box fallback for missing placeholder")
        
        # Add text box in center of slide
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        
        text = content_data.get('text', '')
        
        if isinstance(text, list):
            for item in text:
                p = text_frame.add_paragraph()
                p.text = str(item)
        else:
            text_frame.text = str(text)
    
    def _add_notes(self, slide, notes: str):
        """
        Add speaker notes to slide
        
        Args:
            slide: Slide object
            notes: Notes text
        """
        try:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes
        except Exception as e:
            self.error_handler.log_warning(f"Could not add notes: {e}")
    
    def add_metadata(self, properties: Dict):
        """
        Add document properties/metadata
        
        Args:
            properties: Dictionary of properties
        """
        core_props = self.prs.core_properties
        
        if 'title' in properties:
            core_props.title = properties['title']
        if 'author' in properties:
            core_props.author = properties['author']
        if 'subject' in properties:
            core_props.subject = properties['subject']
        if 'keywords' in properties:
            core_props.keywords = properties['keywords']
        if 'comments' in properties:
            core_props.comments = properties['comments']
