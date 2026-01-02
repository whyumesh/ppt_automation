"""
Template analysis and structure extraction
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from dataclasses import dataclass, asdict
from typing import List, Dict, Optional
from pathlib import Path
import json
import hashlib

from .error_handler import ErrorHandler, TemplateError


@dataclass
class PlaceholderInfo:
    """Information about a placeholder"""
    placeholder_idx: int
    placeholder_type: str
    max_chars: int
    font_name: str
    font_size: float
    width: float
    height: float
    left: float
    top: float


@dataclass
class SlideTemplate:
    """Template for a single slide"""
    slide_idx: int
    layout_name: str
    placeholders: List[PlaceholderInfo]
    slide_type: str
    has_title: bool
    has_body: bool


class TemplateAnalyzer:
    """Analyze PowerPoint templates and extract structure"""
    
    def __init__(self, template_path: str, cache_dir: str = "cache/"):
        self.template_path = Path(template_path)
        self.cache_dir = Path(cache_dir)
        self.cache_dir.mkdir(exist_ok=True)
        
        self.error_handler = ErrorHandler()
        self.prs: Optional[Presentation] = None
        self.structure: Dict = {}
        
        # Validate template exists
        ErrorHandler.validate_file_exists(str(self.template_path))
    
    def analyze(self, use_cache: bool = True) -> Dict:
        """
        Analyze template and return structure
        
        Args:
            use_cache: Whether to use cached analysis
            
        Returns:
            Dictionary containing template structure
        """
        cache_file = self._get_cache_path()
        
        # Try to load from cache
        if use_cache and cache_file.exists():
            self.error_handler.log_info(f"Loading template from cache: {cache_file}")
            with open(cache_file, 'r') as f:
                self.structure = json.load(f)
                return self.structure
        
        # Perform fresh analysis
        self.error_handler.log_info(f"Analyzing template: {self.template_path}")
        self.prs = Presentation(str(self.template_path))
        
        self.structure = {
            'template_path': str(self.template_path),
            'template_hash': self._calculate_hash(),
            'slides': self._analyze_slides(),
            'theme': self._extract_theme(),
            'layouts': self._extract_layouts(),
            'total_slides': len(self.prs.slides)
        }
        
        # Save to cache
        if use_cache:
            with open(cache_file, 'w') as f:
                json.dump(self.structure, f, indent=2)
            self.error_handler.log_info(f"Template analysis cached: {cache_file}")
        
        return self.structure
    
    def _get_cache_path(self) -> Path:
        """Generate cache file path"""
        template_name = self.template_path.stem
        return self.cache_dir / f"{template_name}_structure.json"
    
    def _calculate_hash(self) -> str:
        """Calculate template file hash for cache validation"""
        with open(self.template_path, 'rb') as f:
            return hashlib.md5(f.read()).hexdigest()
    
    def _analyze_slides(self) -> List[Dict]:
        """Analyze all slides in template"""
        slides_info = []
        
        for slide_idx, slide in enumerate(self.prs.slides):
            try:
                slide_info = SlideTemplate(
                    slide_idx=slide_idx,
                    layout_name=slide.slide_layout.name,
                    placeholders=self._extract_placeholders(slide),
                    slide_type=self._classify_slide(slide),
                    has_title=self._has_title_placeholder(slide),
                    has_body=self._has_body_placeholder(slide)
                )
                slides_info.append(asdict(slide_info))
            except Exception as e:
                self.error_handler.log_error(f"Error analyzing slide {slide_idx}", e)
                continue
        
        return slides_info
    
    def _extract_placeholders(self, slide) -> List[Dict]:
        """Extract placeholder information from slide"""
        placeholders = []
        
        for shape in slide.shapes:
            if not shape.is_placeholder:
                continue
            
            try:
                ph = shape.placeholder_format
                
                # Get font info safely
                font_name = "Arial"  # default
                font_size = 18.0  # default
                
                if shape.has_text_frame and shape.text_frame.paragraphs:
                    para = shape.text_frame.paragraphs[0]
                    if para.runs:
                        run = para.runs[0]
                        if run.font.name:
                            font_name = run.font.name
                        if run.font.size:
                            font_size = run.font.size.pt
                
                placeholder_info = PlaceholderInfo(
                    placeholder_idx=ph.idx,
                    placeholder_type=str(ph.type),
                    max_chars=self._estimate_capacity(shape),
                    font_name=font_name,
                    font_size=font_size,
                    width=shape.width.inches,
                    height=shape.height.inches,
                    left=shape.left.inches,
                    top=shape.top.inches
                )
                
                placeholders.append(asdict(placeholder_info))
                
            except Exception as e:
                self.error_handler.log_warning(f"Could not extract placeholder info: {e}")
                continue
        
        return placeholders
    
    def _estimate_capacity(self, shape) -> int:
        """
        Estimate character capacity based on shape dimensions
        
        Heuristic: ~12 chars per inch width, ~8 lines per inch height,
        ~50 chars per line average
        """
        width = shape.width.inches
        height = shape.height.inches
        
        # Approximate calculation
        chars_per_line = int(width * 12)
        lines = int(height * 8)
        
        return chars_per_line * lines
    
    def _classify_slide(self, slide) -> str:
        """Classify slide type based on layout and content"""
        layout_name = slide.slide_layout.name.lower()
        
        if 'title' in layout_name and 'only' in layout_name:
            return 'title'
        elif 'section' in layout_name:
            return 'section_header'
        elif 'two' in layout_name and 'content' in layout_name:
            return 'two_column'
        elif 'blank' in layout_name:
            return 'blank'
        elif 'content' in layout_name or 'bullet' in layout_name:
            return 'content'
        else:
            return 'content'  # Default
    
    def _has_title_placeholder(self, slide) -> bool:
        """Check if slide has title placeholder"""
        for shape in slide.shapes:
            if shape.is_placeholder:
                ph_type = str(shape.placeholder_format.type)
                if 'TITLE' in ph_type or 'CENTERED_TITLE' in ph_type:
                    return True
        return False
    
    def _has_body_placeholder(self, slide) -> bool:
        """Check if slide has body/content placeholder"""
        for shape in slide.shapes:
            if shape.is_placeholder:
                ph_type = str(shape.placeholder_format.type)
                if 'BODY' in ph_type or 'OBJECT' in ph_type:
                    return True
        return False
    
    def _extract_theme(self) -> Dict:
        """Extract theme information"""
        try:
            theme_info = {
                'master_name': self.prs.slide_master.name if hasattr(self.prs, 'slide_master') else 'Unknown',
                'has_master': hasattr(self.prs, 'slide_master')
            }
        except:
            theme_info = {
                'master_name': 'Unknown',
                'has_master': False
            }
        
        return theme_info
    
    def _extract_layouts(self) -> List[str]:
        """Extract available layout names"""
        try:
            return [layout.name for layout in self.prs.slide_layouts]
        except:
            return []
    
    def get_layout_by_type(self, slide_type: str) -> Optional[int]:
        """Get slide index for a specific slide type"""
        for slide_data in self.structure.get('slides', []):
            if slide_data['slide_type'] == slide_type:
                return slide_data['slide_idx']
        return None
